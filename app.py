import io
import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
from router import route
from agents.gemini_agent import ask as gemini_ask
from agents.groq_agent import ask as groq_ask
from agents.mistral_agent import ask as mistral_ask
from session_manager import list_sessions, load_session, save_session, create_session, delete_session, auto_name

IMAGE_TYPES = (".png", ".jpg", ".jpeg", ".webp", ".gif")
DOC_ICONS = {".pptx": "📊", ".pdf": "📄", ".docx": "📝"}


def extract_file_text(file_bytes: bytes, filename: str) -> str:
    if filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[슬라이드 {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    if filename.endswith(".docx"):
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())

    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(file_bytes))
        lines = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text:
                lines.append(f"[페이지 {i}]\n{text.strip()}")
        return "\n".join(lines)

    return ""

load_dotenv()

AGENTS = {
    "gemini": {"label": "Gemini 2.5 Flash", "icon": "🔵", "desc": "검색 · 글쓰기 · 진로"},
    "groq":   {"label": "Groq / Llama 4",   "icon": "🟠", "desc": "보안 · CTF · 리눅스"},
    "mistral":{"label": "Mistral Codestral", "icon": "🟣", "desc": "코드 · 개발 · 리뷰"},
}

AGENT_FN = {
    "gemini":  gemini_ask,
    "groq":    groq_ask,
    "mistral": mistral_ask,
}

st.set_page_config(page_title="OrchestrAI", page_icon="🤖", layout="wide")

# 세션 초기화
if "current_session_id" not in st.session_state:
    session_id = create_session()
    st.session_state.current_session_id = session_id
    st.session_state.messages = []
    st.session_state.session_name = "새 대화"
if "uploader_key" not in st.session_state:
    st.session_state.uploader_key = 0

# 사이드바
with st.sidebar:
    st.title("🤖 OrchestrAI")
    st.caption("질문에 따라 최적의 AI를 자동 선택합니다")
    st.divider()

    if st.button("✏️ 새 대화", use_container_width=True):
        session_id = create_session()
        st.session_state.current_session_id = session_id
        st.session_state.messages = []
        st.session_state.session_name = "새 대화"
        st.session_state.uploader_key += 1
        st.rerun()

    st.divider()
    st.markdown("#### 대화 목록")

    sessions = list_sessions()
    for s in sessions:
        col1, col2 = st.columns([5, 1])
        is_current = s["id"] == st.session_state.current_session_id
        label = f"{'▶ ' if is_current else ''}{s['name']}"
        with col1:
            if st.button(label, key=f"sess_{s['id']}", use_container_width=True):
                save_session(
                    st.session_state.current_session_id,
                    st.session_state.session_name,
                    st.session_state.messages,
                )
                loaded = load_session(s["id"])
                st.session_state.current_session_id = s["id"]
                st.session_state.messages = loaded["messages"]
                st.session_state.session_name = loaded["name"]
                st.session_state.uploader_key += 1
                st.rerun()
        with col2:
            if st.button("🗑", key=f"del_{s['id']}"):
                delete_session(s["id"])
                if is_current:
                    new_id = create_session()
                    st.session_state.current_session_id = new_id
                    st.session_state.messages = []
                    st.session_state.session_name = "새 대화"
                st.rerun()

    st.divider()
    st.markdown("#### 접두어로 AI 직접 지정")
    for prefix, agent in [("gemini:", "gemini"), ("groq:", "groq"), ("code:", "mistral")]:
        info = AGENTS[agent]
        st.markdown(f"`{prefix}` → {info['icon']} {info['label']}")
    st.divider()
    st.markdown("#### 담당 AI")
    for key, info in AGENTS.items():
        st.markdown(f"{info['icon']} {info['label']}  \n{info['desc']}")

# 대화 기록 출력
for msg in st.session_state.messages:
    if msg["role"] == "user":
        with st.chat_message("user"):
            if msg.get("image"):
                st.image(msg["image"], width=300)
            st.write(msg["content"])
    else:
        info = AGENTS[msg.get("agent", "gemini")]
        with st.chat_message("assistant"):
            st.caption(f"{info['icon']} {info['label']}")
            st.write(msg["content"])

# 파일 업로드
uploaded_file = st.file_uploader(
    "파일 첨부 (선택)",
    type=["png", "jpg", "jpeg", "webp", "gif", "pptx", "pdf", "docx"],
    label_visibility="collapsed",
    key=f"uploader_{st.session_state.uploader_key}",
)
if uploaded_file:
    fname = uploaded_file.name.lower()
    if any(fname.endswith(ext) for ext in IMAGE_TYPES):
        st.image(uploaded_file, width=200, caption="전송 시 Gemini로 자동 분석")
    else:
        icon = next((v for k, v in DOC_ICONS.items() if fname.endswith(k)), "📎")
        st.caption(f"{icon} {uploaded_file.name} — 전송 시 텍스트 추출 후 분석")

# 입력
if prompt := st.chat_input("질문을 입력하세요..."):
    fname = uploaded_file.name.lower() if uploaded_file else ""
    is_image = uploaded_file and any(fname.endswith(ext) for ext in IMAGE_TYPES)
    is_doc   = uploaded_file and not is_image

    image_bytes = uploaded_file.getvalue() if is_image else None
    mime_type   = uploaded_file.type if is_image else None

    if is_doc:
        doc_text = extract_file_text(uploaded_file.getvalue(), fname)
        clean_query = f"다음 파일 내용을 참고해서 답해줘:\n\n{doc_text}\n\n질문: {prompt}"
        agent, _ = route(prompt)
    elif is_image:
        agent, clean_query = "gemini", prompt
    else:
        agent, clean_query = route(prompt)
    info = AGENTS[agent]

    history = [
        {"role": m["role"], "content": m["content"]}
        for m in st.session_state.messages[-10:]
    ]

    with st.spinner(f"{info['icon']} {info['label']} 처리 중..."):
        try:
            if image_bytes:
                response = gemini_ask(clean_query, history, image_bytes, mime_type)
            else:
                response = AGENT_FN[agent](clean_query, history)
        except Exception as e:
            err = str(e)
            if "429" in err or "RESOURCE_EXHAUSTED" in err:
                import re
                retry = re.search(r'retry[^0-9]*(\d+)', err)
                wait = f"{retry.group(1)}초 후 재시도" if retry else "잠시 후 재시도"
                response = f"⚠️ {info['label']} 무료 한도 초과 — {wait}\n`groq:` 또는 `gemini:` 접두어로 다른 AI 사용 가능"
            else:
                response = f"오류: {err}"

    st.session_state.messages.append({"role": "user", "content": prompt, "image": image_bytes})
    st.session_state.messages.append({"role": "assistant", "content": response, "agent": agent})

    # 첫 메시지 기준으로 세션 이름 자동 설정
    if len(st.session_state.messages) == 2:
        st.session_state.session_name = auto_name(st.session_state.messages)

    save_session(
        st.session_state.current_session_id,
        st.session_state.session_name,
        st.session_state.messages,
    )
    st.session_state.uploader_key += 1
    st.rerun()
